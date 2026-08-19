# Copyright 2026 徐松夏（Xu Songxia）
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""PowerPoint (.pptx) document parser using python-pptx."""

from pathlib import Path

from pptx import Presentation

from app.parsers.base import BaseParser, ParsedDocument, ParsedSection, ParserPluginMeta


class PptxParser(BaseParser):
    """Parse .pptx files, producing one section per slide."""

    def extensions(self) -> list[str]:
        return ["pptx"]

    @classmethod
    def plugin_meta(cls) -> ParserPluginMeta:
        return ParserPluginMeta(
            name="pptx",
            display_name="parser.pptx.name",
            description="parser.pptx.desc",
            category="office",
            extensions=["pptx"],
        )

    def parse(self, file_path: Path) -> ParsedDocument:
        prs = Presentation(file_path)

        sections: list[ParsedSection] = []
        for idx, slide in enumerate(prs.slides):
            title, body_parts = self._extract_slide_text(slide)
            if not title and not body_parts:
                continue
            content = "\n".join(body_parts) if body_parts else title or ""
            sections.append(ParsedSection(
                level=1,
                heading=title or f"Slide {idx + 1}",
                content=content,
                page=idx + 1,
            ))

        if not sections:
            sections.append(ParsedSection(
                level=0, heading=file_path.stem, content="",
            ))

        return ParsedDocument(
            title=file_path.stem, file_type="pptx", sections=sections,
            metadata={"slide_count": len(prs.slides)},
        )

    def _extract_slide_text(self, slide) -> tuple[str, list[str]]:
        title = ""
        body_parts: list[str] = []

        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()

        for shape in slide.shapes:
            # Skip the title shape (already captured above)
            if shape == slide.shapes.title:
                continue
            text = self._extract_shape_text(shape)
            if text:
                body_parts.append(text)

        return title, body_parts

    def _extract_shape_text(self, shape) -> str:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    rows.append("\t".join(cells))
            return "\n".join(rows)
        return ""
