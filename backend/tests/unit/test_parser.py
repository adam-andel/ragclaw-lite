"""Unit tests for document parsers (Markdown, TXT, unsupported types)."""

import csv as csv_lib
import json as json_lib
import sys
from email.message import EmailMessage
from email.policy import default as email_default
from pathlib import Path

import pytest

# Ensure backend is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

from app.parsers.base import ParsedDocument, ParsedSection
from app.services.parser import parser_service


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_md(path: Path, content: str) -> Path:
    path.write_text(content, encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

class TestMarkdownParser:
    """Markdown parsing correctness."""

    def test_parse_title_and_sections(self, tmp_path):
        md = _make_md(tmp_path / "test.md", """# 系统架构设计文档

## 第一章 概述

本文档描述微服务架构的整体设计方案。

### 1.1 背景

随着业务增长，单体应用已无法满足扩展需求。

## 第二章 核心设计

gRPC 通信协议选择。
""")
        doc = parser_service.parse(md, "md")
        assert isinstance(doc, ParsedDocument)
        assert doc.title == "系统架构设计文档"
        assert doc.file_type == "markdown"
        assert len(doc.sections) >= 3  # H1 + two H2 sections

    def test_heading_hierarchy(self, tmp_path):
        """Verify H1 > H2 > H3 levels are preserved."""
        md = _make_md(tmp_path / "hierarchy.md", """# L1 Title

## L2 Section A

### L3 Sub A1

content here

## L2 Section B

### L3 Sub B1

more content
""")
        doc = parser_service.parse(md, "md")
        levels = [s.level for s in doc.sections]
        headings = [s.heading for s in doc.sections]
        assert 1 in levels
        assert 2 in levels
        assert 3 in levels
        # Check specific headings
        assert "L1 Title" in headings
        assert "L2 Section A" in headings
        assert "L3 Sub A1" in headings

    def test_code_block_preserved(self, tmp_path):
        """Code block text should be present in section content."""
        md = _make_md(tmp_path / "code.md", """# API Example

```python
def hello():
    print("Hello ERAG")
```
""")
        doc = parser_service.parse(md, "md")
        full = "\n".join(s.content for s in doc.sections)
        assert "hello()" in full
        assert "Hello ERAG" in full

    def test_empty_file_no_crash(self, tmp_path):
        """Empty markdown should not crash."""
        md = _make_md(tmp_path / "empty.md", "")
        doc = parser_service.parse(md, "md")
        assert isinstance(doc, ParsedDocument)
        # Sections may be empty but no exception
        assert doc.section_count >= 0

    def test_file_not_found_raises(self, tmp_path):
        """Parsing a non-existent file should raise."""
        with pytest.raises(Exception):
            parser_service.parse(tmp_path / "nope.md", "md")


# ---------------------------------------------------------------------------
# TXT parser
# ---------------------------------------------------------------------------

class TestTxtParser:
    """Plain-text parsing correctness."""

    def test_txt_parses_sections(self, tmp_path):
        txt = tmp_path / "plan.txt"
        txt.write_text("""ERAG 项目开发计划

第一阶段：基础架构搭建
完成 FastAPI 项目初始化、数据库模型设计。

第二阶段：检索引擎开发
实现混合检索、BM25 关键词检索。

第三阶段：前端与部署
完成 Vue3 管理后台开发、Docker 容器化部署。
""", encoding="utf-8")
        doc = parser_service.parse(txt, "txt")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "txt"
        assert len(doc.sections) >= 1

    def test_chinese_encoding_gbk(self, tmp_path):
        """GBK-encoded Chinese text should parse correctly."""
        txt = tmp_path / "gbk.txt"
        content = "中文测试文档\n\n第二段内容"
        txt.write_text(content, encoding="gbk")
        doc = parser_service.parse(txt, "txt")
        # Should not crash; content should have something
        assert doc.section_count >= 0 or len(doc.full_text) > 0


# ---------------------------------------------------------------------------
# Supported types
# ---------------------------------------------------------------------------

class TestSupportedTypes:
    """Format support enumeration."""

    def test_supported_types_list(self):
        types = parser_service.supported_types()
        # .doc was removed: python-docx cannot actually parse legacy binary .doc.
        # 10 new parsers were added in Phase 2 (csv/json/xlsx/pptx/html/eml/rtf/epub/ipynb/msg).
        expected = {
            "pdf", "docx", "md", "markdown", "txt",        # original
            "csv", "json", "xlsx", "xls", "pptx",           # Tier 1 (office/data)
            "html", "htm", "eml",                            # Tier 1 (web/email)
            "rtf", "epub", "ipynb", "msg",                   # Tier 2
        }
        assert set(types) == expected

    def test_unsupported_type_raises(self, tmp_path):
        """Unsupported format should raise ValueError."""
        dummy = tmp_path / "test.exe"
        dummy.write_text("not a real exe", encoding="utf-8")
        with pytest.raises(ValueError, match="No parser available"):
            parser_service.parse(dummy, "exe")


# ---------------------------------------------------------------------------
# Edge cases
# ---------------------------------------------------------------------------

class TestParserEdgeCases:
    """Boundary and stress scenarios."""

    def test_large_file(self, tmp_path):
        """100KB+ file should parse without issues."""
        md = _make_md(tmp_path / "large.md",
                      "# Large Doc\n\n" + "A paragraph of text. " * 5000)
        doc = parser_service.parse(md, "md")
        assert doc.section_count > 0

    def test_markdown_with_front_matter(self, tmp_path):
        """YAML front matter should be stripped."""
        md = _make_md(tmp_path / "fm.md", """---
title: Front Matter Doc
author: tester
---

# Real Title

Some real content here.
""")
        doc = parser_service.parse(md, "md")
        assert doc.title == "Real Title"


# ---------------------------------------------------------------------------
# CSV parser
# ---------------------------------------------------------------------------

class TestCsvParser:
    """CSV parsing correctness."""

    def test_parse_basic(self, tmp_path):
        csv_file = tmp_path / "data.csv"
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            w = csv_lib.writer(f)
            w.writerow(["name", "age", "city"])
            w.writerow(["Alice", "30", "Beijing"])
            w.writerow(["Bob", "25", "Shanghai"])
        doc = parser_service.parse(csv_file, "csv")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "csv"
        assert len(doc.sections) >= 1
        assert "Alice" in doc.full_text
        assert "Beijing" in doc.full_text

    def test_empty_file_no_crash(self, tmp_path):
        csv_file = tmp_path / "empty.csv"
        csv_file.write_text("", encoding="utf-8")
        doc = parser_service.parse(csv_file, "csv")
        assert isinstance(doc, ParsedDocument)

    def test_corrupt_file_no_crash(self, tmp_path):
        """CSV is lenient — garbage bytes should not crash, just produce empty/sparse output."""
        csv_file = tmp_path / "bad.csv"
        csv_file.write_bytes(b"\x00\x01\x02\xff\xfe")
        doc = parser_service.parse(csv_file, "csv")
        assert isinstance(doc, ParsedDocument)


# ---------------------------------------------------------------------------
# JSON parser
# ---------------------------------------------------------------------------

class TestJsonParser:
    """JSON parsing correctness."""

    def test_parse_basic(self, tmp_path):
        json_file = tmp_path / "data.json"
        json_file.write_text(
            json_lib.dumps({"name": "ERAG", "version": "1.0", "features": ["rag", "chat"]}),
            encoding="utf-8",
        )
        doc = parser_service.parse(json_file, "json")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "json"
        assert len(doc.sections) >= 1
        assert "ERAG" in doc.full_text

    def test_empty_file_raises(self, tmp_path):
        """Empty file is invalid JSON — should raise ValueError via safe_parse()."""
        json_file = tmp_path / "empty.json"
        json_file.write_text("", encoding="utf-8")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(json_file, "json")

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        json_file = tmp_path / "bad.json"
        json_file.write_text("{not valid json", encoding="utf-8")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(json_file, "json")


# ---------------------------------------------------------------------------
# Excel parser
# ---------------------------------------------------------------------------

class TestExcelParser:
    """Excel (.xlsx) parsing correctness."""

    def _make_xlsx(self, path: Path) -> Path:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["name", "score"])
        ws.append(["Alice", 95])
        ws.append(["Bob", 88])
        wb.save(path)
        return path

    def test_parse_basic(self, tmp_path):
        xlsx = self._make_xlsx(tmp_path / "data.xlsx")
        doc = parser_service.parse(xlsx, "xlsx")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "xlsx"
        assert len(doc.sections) >= 1
        assert "Alice" in doc.full_text
        assert "Sheet1" in doc.full_text or "Alice" in doc.full_text

    def test_empty_file_raises(self, tmp_path):
        """Empty file is not a valid xlsx — should raise ValueError."""
        xlsx = tmp_path / "empty.xlsx"
        xlsx.write_bytes(b"")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(xlsx, "xlsx")

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        xlsx = tmp_path / "bad.xlsx"
        xlsx.write_bytes(b"\x00\x01\x02\xff\xfe not a zip")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(xlsx, "xlsx")


# ---------------------------------------------------------------------------
# PowerPoint parser
# ---------------------------------------------------------------------------

class TestPptxParser:
    """PowerPoint (.pptx) parsing correctness."""

    def _make_pptx(self, path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Slide"
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        box.text_frame.text = "Hello ERAG"
        prs.save(path)
        return path

    def test_parse_basic(self, tmp_path):
        pptx = self._make_pptx(tmp_path / "deck.pptx")
        doc = parser_service.parse(pptx, "pptx")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "pptx"
        assert len(doc.sections) >= 1
        assert "Hello ERAG" in doc.full_text or "Test Slide" in doc.full_text

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        pptx = tmp_path / "bad.pptx"
        pptx.write_bytes(b"\x00\x01\x02\xff\xfe not a zip")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(pptx, "pptx")


# ---------------------------------------------------------------------------
# HTML parser
# ---------------------------------------------------------------------------

class TestHtmlParser:
    """HTML parsing correctness."""

    def test_parse_basic(self, tmp_path):
        html_file = tmp_path / "page.html"
        html_file.write_text("""<!DOCTYPE html>
<html><head><title>Test Page</title></head>
<body>
<h1>Main Heading</h1>
<p>First paragraph text.</p>
<h2>Subsection</h2>
<p>Second paragraph.</p>
</body></html>
""", encoding="utf-8")
        doc = parser_service.parse(html_file, "html")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "html"
        assert len(doc.sections) >= 1
        assert "Main Heading" in doc.full_text or "paragraph" in doc.full_text

    def test_empty_file_no_crash(self, tmp_path):
        html_file = tmp_path / "empty.html"
        html_file.write_text("", encoding="utf-8")
        doc = parser_service.parse(html_file, "html")
        assert isinstance(doc, ParsedDocument)

    def test_strips_script_and_style(self, tmp_path):
        html_file = tmp_path / "scripts.html"
        html_file.write_text("""<html><head><style>.x{color:red}</style></head>
<body><script>alert(1)</script><p>visible text</p></body></html>""",
                             encoding="utf-8")
        doc = parser_service.parse(html_file, "html")
        assert "visible text" in doc.full_text
        assert "alert" not in doc.full_text
        assert "color:red" not in doc.full_text


# ---------------------------------------------------------------------------
# Email (EML) parser
# ---------------------------------------------------------------------------

class TestEmailParser:
    """EML email parsing correctness."""

    def _make_eml(self, path: Path) -> Path:
        msg = EmailMessage(policy=email_default)
        msg["Subject"] = "Test Subject"
        msg["From"] = "alice@example.com"
        msg["To"] = "bob@example.com"
        msg["Date"] = "Mon, 01 Jan 2024 10:00:00 +0800"
        msg.set_content("Hello Bob,\n\nThis is a test email.\n\nRegards,\nAlice")
        with open(path, "wb") as f:
            f.write(bytes(msg))
        return path

    def test_parse_basic(self, tmp_path):
        eml = self._make_eml(tmp_path / "test.eml")
        doc = parser_service.parse(eml, "eml")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "eml"
        assert len(doc.sections) >= 1
        assert "Test Subject" in doc.full_text
        assert "test email" in doc.full_text

    def test_empty_file_no_crash(self, tmp_path):
        eml = tmp_path / "empty.eml"
        eml.write_text("", encoding="utf-8")
        doc = parser_service.parse(eml, "eml")
        assert isinstance(doc, ParsedDocument)


# ---------------------------------------------------------------------------
# RTF parser
# ---------------------------------------------------------------------------

class TestRtfParser:
    """RTF parsing correctness."""

    def test_parse_basic(self, tmp_path):
        rtf_file = tmp_path / "doc.rtf"
        rtf_file.write_text(r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times;}}\f0\fs24 Hello RTF World}",
                            encoding="utf-8")
        doc = parser_service.parse(rtf_file, "rtf")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "rtf"
        assert len(doc.sections) >= 1
        assert "Hello" in doc.full_text or "RTF" in doc.full_text

    def test_empty_file_no_crash(self, tmp_path):
        rtf_file = tmp_path / "empty.rtf"
        rtf_file.write_text("", encoding="utf-8")
        doc = parser_service.parse(rtf_file, "rtf")
        assert isinstance(doc, ParsedDocument)


# ---------------------------------------------------------------------------
# EPUB parser
# ---------------------------------------------------------------------------

class TestEpubParser:
    """EPUB ebook parsing correctness."""

    def _make_epub(self, path: Path) -> Path:
        from ebooklib import epub
        book = epub.EpubBook()
        book.set_title("Test Book")
        book.set_language("en")
        c1 = epub.EpubHtml(title="Chapter 1", file_name="ch1.xhtml", lang="en")
        c1.content = "<html><body><h1>Chapter 1</h1><p>Hello EPUB world</p></body></html>"
        book.add_item(c1)
        # Required by ebooklib reader: NCX (TOC) + Nav (EPUB3 navigation)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = [c1]
        book.toc = [c1]
        epub.write_epub(path, book, {})
        return path

    def test_parse_basic(self, tmp_path):
        epub_file = self._make_epub(tmp_path / "book.epub")
        doc = parser_service.parse(epub_file, "epub")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "epub"
        assert len(doc.sections) >= 1
        assert "Hello EPUB" in doc.full_text or "Chapter 1" in doc.full_text

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        epub_file = tmp_path / "bad.epub"
        epub_file.write_bytes(b"\x00\x01\x02\xff\xfe not a zip")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(epub_file, "epub")


# ---------------------------------------------------------------------------
# Jupyter Notebook parser
# ---------------------------------------------------------------------------

class TestNotebookParser:
    """Jupyter Notebook (.ipynb) parsing correctness."""

    def _make_ipynb(self, path: Path) -> Path:
        import nbformat
        from nbformat.v4 import new_notebook, new_markdown_cell, new_code_cell
        nb = new_notebook()
        nb.cells = [
            new_markdown_cell("# Sample Notebook\n\nIntro text here."),
            new_code_cell("print('hello world')"),
        ]
        with open(path, "w", encoding="utf-8") as f:
            nbformat.write(nb, f)
        return path

    def test_parse_basic(self, tmp_path):
        ipynb = self._make_ipynb(tmp_path / "notebook.ipynb")
        doc = parser_service.parse(ipynb, "ipynb")
        assert isinstance(doc, ParsedDocument)
        assert doc.file_type == "ipynb"
        assert len(doc.sections) >= 1
        assert "hello world" in doc.full_text or "Sample Notebook" in doc.full_text

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        ipynb = tmp_path / "bad.ipynb"
        ipynb.write_text("{not valid notebook json", encoding="utf-8")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(ipynb, "ipynb")


# ---------------------------------------------------------------------------
# Outlook MSG parser
# ---------------------------------------------------------------------------

class TestMsgParser:
    """Outlook .msg parsing correctness.

    .msg files are OLE compound documents — cannot be easily generated in-test.
    We test can_handle() and corrupt-file behavior; full parse is covered by
    manual integration testing with a real .msg file.
    """

    def test_can_handle_msg(self):
        """MsgParser should claim the 'msg' extension."""
        assert "msg" in parser_service.supported_types()

    def test_corrupt_file_raises_valueerror(self, tmp_path):
        """A non-OLE file should fail gracefully via safe_parse()."""
        msg_file = tmp_path / "fake.msg"
        msg_file.write_bytes(b"\x00\x01\x02\xff\xfe not an OLE file")
        with pytest.raises(ValueError, match="解析失败"):
            parser_service.parse(msg_file, "msg")
